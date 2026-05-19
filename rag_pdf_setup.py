"""
Enhanced RAG Setup: Unified Vector Store from PDFs + Knowledge Base + PPT
--------------------------------------------------------------------------
Creates a FAISS vector database from:
1. knowledge.txt (comprehensive healthcare knowledge base)
2. documents/*.pdf (government reports)
3. data/*.pptx (presentation data)
"""

import os
from langchain_community.document_loaders import TextLoader, PyPDFLoader
from langchain_community.vectorstores import FAISS
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document

print("=" * 60)
print("Unified RAG Vector Store Builder")
print("=" * 60)

all_documents = []

# -------------------------
# 1. Load Knowledge Base
# -------------------------
print("\n[1/4] Loading knowledge base (knowledge.txt)...")
loader = TextLoader("knowledge.txt", encoding="utf-8")
kb_docs = loader.load()
all_documents.extend(kb_docs)
print(f"   Loaded {len(kb_docs)} document(s) from knowledge.txt")

# -------------------------
# 2. Load PDF Documents
# -------------------------
print("\n[2/4] Loading PDF documents from documents/...")
pdf_count = 0
if os.path.isdir("documents"):
    for file in os.listdir("documents"):
        if file.endswith(".pdf"):
            try:
                loader = PyPDFLoader(f"documents/{file}")
                pdf_docs = loader.load()
                all_documents.extend(pdf_docs)
                pdf_count += 1
                print(f"   [OK] {file} ({len(pdf_docs)} pages)")
            except Exception as e:
                print(f"   [FAIL] {file}: {e}")
else:
    print("   documents/ not found, skipping PDF documents")
print(f"   Total PDFs processed: {pdf_count}")

# -------------------------
# 3. Load PPT Content
# -------------------------
print("\n[3/4] Loading PowerPoint presentations from data/...")
ppt_count = 0
try:
    from pptx import Presentation

    for file in os.listdir("data"):
        if file.endswith(".pptx"):
            try:
                prs = Presentation(f"data/{file}")
                ppt_text = []
                for i, slide in enumerate(prs.slides):
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text.strip())
                    if slide_text:
                        ppt_text.append(f"Slide {i+1}: " + " | ".join(slide_text))

                if ppt_text:
                    full_text = "\n\n".join(ppt_text)
                    doc = Document(
                        page_content=full_text,
                        metadata={"source": f"data/{file}", "type": "presentation"}
                    )
                    all_documents.append(doc)
                    ppt_count += 1
                    print(f"   [OK] {file} ({len(ppt_text)} slides with content)")
            except Exception as e:
                print(f"   [FAIL] {file}: {e}")
except ImportError:
    print("   python-pptx not installed, skipping PPT files")
print(f"   Total PPTs processed: {ppt_count}")

# -------------------------
# 4. Load CSV Data as Context
# -------------------------
print("\n[3.5/4] Loading CSV data summary...")
try:
    import pandas as pd
    df = pd.read_csv("data/india_healthcare_data.csv")
    # Create summary documents per state (latest year only)
    latest_year = df['year'].max()
    for _, row in df[df['year'] == latest_year].iterrows():
        state_summary = (
            f"State: {row['state']} (Year {latest_year})\n"
            f"Population: {row['population_crore']} crore, Urban: {row['urban_pct']}%\n"
            f"Hospitals: {row['hospitals_total']}, Hospital Beds per 1000: {row['hospital_beds_per_1000']}\n"
            f"ICU Beds: {row['icu_beds']}, Doctors: {row['doctors_total']}, Doctor per 1000: {row['doctor_per_1000']}\n"
            f"Nurses: {row['nurses_total']}, PHCs: {row.get('phc_count', 'N/A')}, CHCs: {row.get('chc_count', 'N/A')}\n"
            f"Vaccine Coverage: {row['vaccine_coverage_pct']}%, Cold Chain Facilities: {row['cold_chain_facilities']}\n"
            f"Health Budget: Rs {row['health_budget_crore']} crore, Per Capita: Rs {row['budget_per_capita_inr']}\n"
            f"Disease Index: {row['disease_index']}, Infrastructure Gap Score: {row['infra_gap_score']}\n"
            f"Maternal Mortality Ratio: {row['maternal_mortality_ratio']}, Infant Mortality Rate: {row['infant_mortality_rate']}\n"
            f"Life Expectancy: {row['life_expectancy']} years"
        )
        doc = Document(
            page_content=state_summary,
            metadata={"source": "india_healthcare_data.csv", "state": row['state'], "year": latest_year}
        )
        all_documents.append(doc)
    print(f"   Created {len(df[df['year'] == latest_year])} state profile documents")
except Exception as e:
    print(f"   Could not load CSV: {e}")

# -------------------------
# 4. Split and Embed
# -------------------------
print(f"\n[4/4] Splitting {len(all_documents)} documents into chunks...")

text_splitter = RecursiveCharacterTextSplitter(
    chunk_size=800,
    chunk_overlap=150,
    separators=["\n\n", "\n", "---", ". ", " "]
)

docs = text_splitter.split_documents(all_documents)
print(f"   Created {len(docs)} chunks")

print("\nCreating embeddings (this may take a few minutes)...")
embeddings = HuggingFaceEmbeddings(
    model_name="sentence-transformers/all-MiniLM-L6-v2"
)

db = FAISS.from_documents(docs, embeddings)

db.save_local("vector_store")

print(f"\n{'='*60}")
print(f"Vector store created successfully!")
print(f"   Total documents: {len(all_documents)}")
print(f"   Total chunks: {len(docs)}")
print(f"   Embedding model: all-MiniLM-L6-v2")
print(f"   Saved to: vector_store/")
print(f"{'='*60}")
